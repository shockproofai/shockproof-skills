#!/usr/bin/env python3
"""Extract text and shape metadata from a PPTX for the text-edit-document workflow.

Produces two JSON files:
  - {deckname}_intermediate.json  (text content for docx generation)
  - {deckname}_extract.json       (shape ID sidecar for round-trip)
"""

import json
import sys
import os
import subprocess
from datetime import datetime, timezone

# ── Auto-install python-pptx into a local vendor/ directory ──────────────
SKILL_DIR = os.path.dirname(os.path.abspath(__file__))
VENDOR_DIR = os.path.join(SKILL_DIR, "vendor")

if not os.path.isdir(VENDOR_DIR):
    print("First run: installing python-pptx into vendor/ …", file=sys.stderr)
    subprocess.check_call(
        [sys.executable, "-m", "pip", "install",
         "--target", VENDOR_DIR, "--quiet", "python-pptx"],
    )
    print("Done.", file=sys.stderr)

# Prepend vendor/ so imports resolve from there
if VENDOR_DIR not in sys.path:
    sys.path.insert(0, VENDOR_DIR)
# ─────────────────────────────────────────────────────────────────────────

from pptx import Presentation
from pptx.util import Emu

# Placeholder type indices for title/subtitle
TITLE_TYPES = {0, 15}  # TITLE, CENTER_TITLE
SUBTITLE_TYPE = 1       # SUBTITLE


def reading_order_key(shape):
    """Sort shapes top-to-bottom, then left-to-right."""
    return (shape.top or 0, shape.left or 0)


def extract_paragraphs(text_frame):
    """Extract paragraph objects from a text frame, preserving bullet/number info.

    Returns list of dicts: {"text": str, "bullet": str|None, "autonum": str|None}
    - bullet: the buChar character (e.g. "•") if present
    - autonum: the buAutoNum type (e.g. "arabicPeriod") if present
    """
    from pptx.oxml.ns import qn
    result = []
    for para in text_frame.paragraphs:
        text = para.text
        bullet = None
        autonum = None
        pPr = para._p.find(qn('a:pPr'))
        if pPr is not None:
            buChar = pPr.find(qn('a:buChar'))
            if buChar is not None:
                bullet = buChar.get('char')
            buAuto = pPr.find(qn('a:buAutoNum'))
            if buAuto is not None:
                autonum = buAuto.get('type')
        result.append({"text": text, "bullet": bullet, "autonum": autonum})
    return result


def classify_shapes(slide):
    """Classify slide shapes into title, subtitle, and body categories.

    Strategy:
    1. First try standard placeholders (TITLE, CENTER_TITLE, SUBTITLE).
    2. If no placeholders found (common with pptxgenjs-generated decks),
       use a heuristic: sort text shapes by vertical position and treat
       the top two as title and subtitle if they are near the top of the
       slide (top < 40% of slide height ~= 2740000 EMU for 10" slides).
    """
    SLIDE_HEIGHT_40PCT = 2740000  # ~40% of standard 6858000 EMU (7.5") height

    title_shape = None
    subtitle_shape = None
    body_shapes = []
    text_shapes = []  # all shapes with non-empty text

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        # Try placeholder-based classification first
        if shape.is_placeholder:
            ph = shape.placeholder_format
            idx = ph.idx
            if idx in TITLE_TYPES and title_shape is None:
                title_shape = shape
                continue
            if idx == SUBTITLE_TYPE and subtitle_shape is None:
                subtitle_shape = shape
                continue

        text_shapes.append(shape)

    # If placeholders found title/subtitle, remaining go to body
    if title_shape or subtitle_shape:
        body_shapes = text_shapes
    else:
        # Heuristic: sort by vertical position, pick top shapes as title/subtitle
        text_shapes.sort(key=reading_order_key)
        remaining = []
        for shape in text_shapes:
            top = shape.top or 0
            if title_shape is None and top < SLIDE_HEIGHT_40PCT:
                title_shape = shape
            elif subtitle_shape is None and top < SLIDE_HEIGHT_40PCT:
                subtitle_shape = shape
            else:
                remaining.append(shape)
        body_shapes = remaining

    # Sort body shapes by visual reading order
    body_shapes.sort(key=reading_order_key)
    return title_shape, subtitle_shape, body_shapes


def detect_native_tables(slide):
    """Extract native PPTX table shapes (real rows/cells/borders).

    These always pass through as tables regardless of column count.
    Returns a list of table dicts (one per table shape found), or None.
    """
    tables = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        num_cols = len(tbl.columns)
        result_rows = []
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                paras = extract_paragraphs(cell.text_frame)
                cells.append([p for p in paras if p["text"].strip()])
            result_rows.append({"cells": cells})

        tables.append({
            "columns": num_cols,
            "rows": result_rows,
            "trailing": [],
            "native": True
        })
    return tables if tables else None


def detect_table_layout(body_shapes):
    """Detect if body shapes form a tabular layout (rows x columns).

    A table layout has:
    - Multiple rows at distinct vertical positions
    - Each row has the same number of shapes at consistent column positions
    - At least 2 columns and 2 rows

    Returns dict {"columns": int, "rows": [{"cells": [paragraphs...]}], "trailing": [paragraphs]}
    or None if not a table layout.
    """
    if len(body_shapes) < 4:
        return None

    ROW_GAP = 640000  # ~0.7 inch — tighter than card detection for table rows
    COL_SNAP = 914400  # 1.0 inch — shapes with left edges this close are same column

    # Step 1: cluster into rows
    sorted_shapes = sorted(body_shapes, key=lambda s: (s.top or 0))
    rows = []
    current_row = [sorted_shapes[0]]
    for shape in sorted_shapes[1:]:
        if (shape.top or 0) - (current_row[0].top or 0) <= ROW_GAP:
            current_row.append(shape)
        else:
            rows.append(current_row)
            current_row = [shape]
    rows.append(current_row)

    if len(rows) < 2:
        return None

    # Step 2: detect column positions from the row(s) with the most shapes
    max_cols = max(len(r) for r in rows)
    if max_cols < 2 or max_cols > 3:
        return None  # tables with 4+ columns are better rendered as cards

    # Find column anchors from all rows with max_cols shapes
    full_rows = [r for r in rows if len(r) == max_cols]
    col_lefts = []  # list of left positions per column index
    for row in full_rows:
        row_sorted = sorted(row, key=lambda s: (s.left or 0))
        for ci, shape in enumerate(row_sorted):
            if ci >= len(col_lefts):
                col_lefts.append([])
            col_lefts[ci].append(shape.left or 0)

    col_anchors = [sum(positions) / len(positions) for positions in col_lefts]

    # Step 3: check that most rows match the column count
    # Allow trailing shapes (e.g., a footer paragraph) that don't fit the pattern
    table_rows = []
    trailing_shapes = []
    for row in rows:
        if len(row) == max_cols:
            row_sorted = sorted(row, key=lambda s: (s.left or 0))
            # Verify shapes align with column anchors
            aligned = True
            for ci, shape in enumerate(row_sorted):
                if ci < len(col_anchors) and abs((shape.left or 0) - col_anchors[ci]) > COL_SNAP:
                    aligned = False
                    break
            if aligned:
                table_rows.append(row_sorted)
            else:
                trailing_shapes.extend(row)
        else:
            trailing_shapes.extend(row)

    if len(table_rows) < 2:
        return None

    # Step 4: build the table data
    result_rows = []
    for row_shapes in table_rows:
        cells = []
        for shape in row_shapes:
            paras = extract_paragraphs(shape.text_frame)
            cell_text = []
            for p in paras:
                if p["text"].strip():
                    cell_text.append(p)
            cells.append(cell_text)
        result_rows.append({"cells": cells})

    trailing = []
    for shape in trailing_shapes:
        for p in extract_paragraphs(shape.text_frame):
            if p["text"].strip():
                trailing.append(p)

    return {
        "columns": max_cols,
        "rows": result_rows,
        "trailing": trailing
    }


def group_into_cards(body_shapes):
    """Group body shapes into card clusters based on spatial proximity.

    A "card" is a set of shapes (icon/label, title, description) that are
    visually grouped on the slide. Detection heuristic:
    1. Cluster shapes into rows by vertical proximity (top within ROW_GAP).
    2. Within each row, sub-cluster by horizontal position (left within COL_GAP).
    3. If this produces 3+ cards with 2-4 shapes each, treat as card layout.
       Otherwise return None (not a card layout — use flat body text).

    Returns list of cards, each card is a list of shapes sorted top-to-bottom.
    """
    if len(body_shapes) < 3:
        return None

    ROW_GAP = 914400  # 1.0 inch in EMU — shapes within this vertical band are same row

    # Step 1: cluster into rows by vertical proximity
    sorted_shapes = sorted(body_shapes, key=lambda s: (s.top or 0))
    rows = []
    current_row = [sorted_shapes[0]]
    for shape in sorted_shapes[1:]:
        if (shape.top or 0) - (current_row[0].top or 0) <= ROW_GAP:
            current_row.append(shape)
        else:
            rows.append(current_row)
            current_row = [shape]
    rows.append(current_row)

    # Step 2: group shapes into cards by clustering their left-edge positions.
    # Shapes belonging to the same card share a similar horizontal position
    # (their left edges are within CARD_PROXIMITY of each other).
    # This works for both side-by-side and vertically-stacked card layouts.

    CARD_PROXIMITY = 1371600  # 1.5 inches — shapes with left edges this close are same card

    # Collect all left positions across all body shapes
    all_shapes = [s for row in rows for s in row]
    all_shapes_sorted = sorted(all_shapes, key=lambda s: (s.left or 0))

    # Cluster by left position to find column anchors
    col_groups = [[all_shapes_sorted[0]]]
    for shape in all_shapes_sorted[1:]:
        # Check if this shape's left is close to the average left of the current group
        group_avg_left = sum(s.left or 0 for s in col_groups[-1]) / len(col_groups[-1])
        if abs((shape.left or 0) - group_avg_left) <= CARD_PROXIMITY:
            col_groups[-1].append(shape)
        else:
            col_groups.append([shape])

    if len(col_groups) < 2:
        return None  # single column, not a card layout

    # Now intersect column groups with rows to produce cards
    # Compute column boundaries as midpoints between group averages
    col_avgs = sorted([
        sum(s.left or 0 for s in g) / len(g) for g in col_groups
    ])
    boundaries = [(col_avgs[i] + col_avgs[i + 1]) / 2 for i in range(len(col_avgs) - 1)]

    def col_index(shape):
        left = shape.left or 0
        for i, b in enumerate(boundaries):
            if left < b:
                return i
        return len(boundaries)

    cards = []
    for row in rows:
        col_buckets = {}
        for shape in row:
            ci = col_index(shape)
            col_buckets.setdefault(ci, []).append(shape)
        for ci in sorted(col_buckets.keys()):
            col_buckets[ci].sort(key=reading_order_key)
            cards.append(col_buckets[ci])

    # Remove empty cards
    cards = [c for c in cards if c]

    # Step 3: validate — require 3+ cards, each with 2+ shapes
    if len(cards) >= 3 and all(len(c) >= 2 for c in cards):
        return cards
    return None


def extract_slide(slide, slide_index):
    """Extract content and metadata from a single slide."""
    title_shape, subtitle_shape, body_shapes = classify_shapes(slide)

    # Build intermediate content
    title_text = None
    if title_shape:
        title_text = "\n".join(p["text"] for p in extract_paragraphs(title_shape.text_frame))

    subtitle_text = None
    if subtitle_shape:
        subtitle_text = "\n".join(p["text"] for p in extract_paragraphs(subtitle_shape.text_frame))

    # Try to detect structured layouts: native tables first, then spatial
    # tables, then cards.
    # Native PPTX tables (real table shapes) — any column count is fine.
    native_tables = detect_native_tables(slide)
    # Use the first native table if found (most slides have at most one).
    table_data = native_tables[0] if native_tables else None
    # Fall back to spatial table detection (capped at 3 columns).
    if not table_data:
        table_data = detect_table_layout(body_shapes)
    cards_data = None
    if not table_data:
        cards = group_into_cards(body_shapes)
        if cards:
            cards_data = []
            for card_shapes in cards:
                card_paras = []
                for shape in card_shapes:
                    for p in extract_paragraphs(shape.text_frame):
                        if p["text"].strip():
                            card_paras.append(p)
                cards_data.append({"paragraphs": card_paras})

    body_data = []
    body_sidecar = []
    for order, shape in enumerate(body_shapes):
        paras = extract_paragraphs(shape.text_frame)
        body_data.append({"order": order, "paragraphs": paras})
        body_sidecar.append({
            "shape_id": shape.shape_id,
            "name": shape.name,
            "order": order
        })

    # Speaker notes (plain text only — no bullets in notes)
    notes_paras = []
    has_notes = False
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            has_notes = True
            notes_paras = [p["text"] for p in extract_paragraphs(slide.notes_slide.notes_text_frame)]

    intermediate = {
        "slide_number": slide_index + 1,
        "title": title_text,
        "subtitle": subtitle_text,
        "body_shapes": body_data,
        "table": table_data,
        "cards": cards_data,
        "speaker_notes": notes_paras
    }

    sidecar = {
        "slide_index": slide_index,
        "slide_number": slide_index + 1,
        "title_shape_id": title_shape.shape_id if title_shape else None,
        "subtitle_shape_id": subtitle_shape.shape_id if subtitle_shape else None,
        "body_shapes": body_sidecar,
        "has_notes": has_notes
    }

    return intermediate, sidecar


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_text.py <input.pptx> [output_dir]", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.dirname(pptx_path)

    deck_name = os.path.splitext(os.path.basename(pptx_path))[0]
    prs = Presentation(pptx_path)

    intermediate_slides = []
    sidecar_slides = []

    for i, slide in enumerate(prs.slides):
        inter, side = extract_slide(slide, i)
        intermediate_slides.append(inter)
        sidecar_slides.append(side)

    intermediate = {
        "source_pptx": os.path.basename(pptx_path),
        "deck_name": deck_name,
        "slide_count": len(prs.slides),
        "slides": intermediate_slides
    }

    sidecar = {
        "source_pptx": os.path.basename(pptx_path),
        "extracted_at": datetime.now(timezone.utc).isoformat(),
        "slides": sidecar_slides
    }

    intermediate_path = os.path.join(output_dir, f"{deck_name}_intermediate.json")
    sidecar_path = os.path.join(output_dir, f"{deck_name}_extract.json")

    with open(intermediate_path, "w") as f:
        json.dump(intermediate, f, indent=2)
    with open(sidecar_path, "w") as f:
        json.dump(sidecar, f, indent=2)

    print(f"Intermediate: {intermediate_path}")
    print(f"Sidecar:      {sidecar_path}")


if __name__ == "__main__":
    main()
